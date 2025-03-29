from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def pdf_to_images(pdf_path):
    images = []
    pages = convert_from_path(pdf_path, 500)  # You can adjust the DPI (here: 500) as needed

    for page_num, page_image in enumerate(pages):
        img_path = f"page_{page_num + 1}.png"
        page_image.save(img_path, 'PNG')
        images.append(img_path)

    return images

def create_pptx(images, output_pptx):
    prs = Presentation()

    for image_path in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use the blank slide layout
        left = top = Inches(1)
        pic = slide.shapes.add_picture(image_path, left, top)

    prs.save(output_pptx)

if __name__ == "__main__":

    pdf_path = "/Users/Tal/Downloads/book_sarit.pdf"  # Replace with the path to your PDF file
    output_pptx = "/Users/Tal/Downloads/output.pptx"  # Replace with the desired output path for the PowerPoint file

    images = pdf_to_images(pdf_path)
    create_pptx(images, output_pptx)
